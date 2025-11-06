#!/usr/bin/env python3
"""
Servidor MCP para EstudIA usando FastMCP
Sistema similar a NotebookLM para gestión de documentos educativos por aula
"""

import asyncio
import json
import io
from typing import Dict, Any, List, Optional
from datetime import datetime

from fastmcp import FastMCP
from pydantic import BaseModel, Field
from PyPDF2 import PdfReader

# Importar nuestros módulos
from .config import config
from .gemini import gemini_client
from .supabase_client import supabase_client

# Crear instancia del servidor FastMCP
mcp = FastMCP("EstudIA MCP Server", version="2.0.0")

# ====== MODELOS DE DATOS ======

class ChatRequest(BaseModel):
    """Modelo para solicitud de chat"""
    message: str = Field(..., description="Mensaje del usuario")
    classroom_id: str = Field(..., description="ID del aula/classroom")
    user_id: Optional[str] = Field(None, description="ID del usuario para mantener contexto")
    session_id: Optional[str] = Field(None, description="ID de sesión para el chat")

# ====== HERRAMIENTAS MCP ======

# ====== FUNCIÓN AUXILIAR PARA generate_embedding ======

async def _generate_embedding_impl(text: str) -> Dict[str, Any]:
    """
    Implementación interna para generar embeddings.
    Esta función es llamada por el tool y también usada internamente.
    """
    if not text or not text.strip():
        return {
            "success": False,
            "error": "El texto no puede estar vacío"
        }
    
    try:
        # Generar embedding usando el cliente existente
        embedding = await gemini_client.generate_embedding(text)
        actual_dim = len(embedding)
        
        return {
            "success": True,
            "embedding": embedding,
            "dimension": actual_dim,
            "text_length": len(text),
            "model": config.GEMINI_EMBED_MODEL,
            "text_preview": text[:100] + ("..." if len(text) > 100 else "")
        }
    
    except Exception as e:
        error_details = str(e)
        return {
            "success": False,
            "error": f"Error generando embedding: {error_details}"
        }


@mcp.tool()
async def generate_embedding(text: str) -> Dict[str, Any]:
    """
    Genera un embedding vector a partir de texto usando Google Gemini.
    
    Convierte texto en un vector numérico de alta dimensionalidad
    que captura el significado semántico del contenido.
    
    Args:
        text: El texto para convertir en embedding
        
    Returns:
        Dict con el embedding generado, dimensiones y metadata
    """
    print(f"\n{'='*60}")
    print("🎯 TOOL: generate_embedding")
    print(f"{'='*60}")
    print(f"📥 Input: {len(text) if text else 0} caracteres")
    
    if not text or not text.strip():
        print("❌ Validación fallida: texto vacío")
        return {
            "success": False,
            "error": "El texto no puede estar vacío"
        }
    
    print(f"🔄 Generando embedding con Gemini...")
    print(f"   📐 Modelo: {config.GEMINI_EMBED_MODEL}")
    print(f"   📊 Dimensiones: {config.EMBED_DIM}")
    
    result = await _generate_embedding_impl(text)
    
    if result.get("success"):
        print(f"✅ Embedding generado exitosamente ({result.get('dimension')} dims)")
    else:
        print(f"❌ ERROR generando embedding: {result.get('error')}")
    
    return result


# ====== FUNCIÓN AUXILIAR PARA extract_text_from_image ======

async def _extract_text_from_image_impl(
    storage_path: str,
    bucket_name: str = "uploads"
) -> Dict[str, Any]:
    """
    Implementación interna para extraer texto de una imagen usando OCR de Gemini Vision.
    Esta función es llamada internamente y NO es un tool MCP expuesto.
    
    Soporta fotos de apuntes, documentos escaneados, capturas de pantalla, etc.
    Útil para procesar contenido educativo que está en formato de imagen.
    
    Args:
        storage_path: Ruta del archivo de imagen en Supabase Storage
        bucket_name: Nombre del bucket donde está la imagen (default: "uploads")
        
    Returns:
        Dict con el texto extraído y metadata
    """
    print(f"\n{'='*60}")
    print("🔧 INTERNAL: _extract_text_from_image_impl")
    print(f"{'='*60}")
    print(f"📥 Parámetros:")
    print(f"   - storage_path: {storage_path}")
    print(f"   - bucket: {bucket_name}")
    
    try:
        # Paso 1: Descargar la imagen desde Supabase Storage
        print("   🔄 PASO 1: Descargando imagen desde Storage...")
        
        image_data = await asyncio.to_thread(
            lambda: supabase_client.client.storage.from_(bucket_name).download(storage_path)
        )
        
        if not image_data:
            raise Exception("No se pudo descargar la imagen del Storage")
        
        print(f"   ✅ Imagen descargada ({len(image_data)} bytes)")
        
        # Paso 2: Detectar tipo MIME de la imagen
        file_extension = storage_path.lower().split('.')[-1]
        mime_types = {
            'jpg': 'image/jpeg',
            'jpeg': 'image/jpeg',
            'png': 'image/png',
            'gif': 'image/gif',
            'webp': 'image/webp',
            'bmp': 'image/bmp'
        }
        mime_type = mime_types.get(file_extension, 'image/jpeg')
        print(f"   📸 Tipo de imagen: {mime_type}")
        
        # Paso 3: Extraer texto usando Gemini Vision OCR
        print("   🔄 PASO 2: Extrayendo texto con Gemini Vision OCR...")
        extracted_text = await gemini_client.extract_text_from_image(image_data, mime_type)
        
        print(f"   ✅ Texto extraído ({len(extracted_text)} caracteres)")
        print(f"   📄 Preview: {extracted_text[:100]}...")
        
        print(f"✅ OCR completado exitosamente")
        print(f"{'='*60}\n")
        
        return {
            "success": True,
            "extracted_text": extracted_text,
            "text_length": len(extracted_text),
            "storage_path": storage_path,
            "bucket": bucket_name,
            "mime_type": mime_type,
            "preview": extracted_text[:200] + ("..." if len(extracted_text) > 200 else "")
        }
        
    except Exception as e:
        error_details = str(e)
        print(f"\n❌ ERROR en OCR: {error_details}")
        
        return {
            "success": False,
            "error": f"Error extrayendo texto de imagen: {error_details}"
        }


# ====== FUNCIÓN AUXILIAR PARA store_document_chunk ======

async def _store_document_chunk_impl(
    classroom_document_id: str,
    chunk_index: int,
    content: str,
    token_count: Optional[int] = None
) -> Dict[str, Any]:
    """
    Implementación interna para almacenar un chunk individual.
    Esta función es llamada por el tool y también usada internamente.
    """
    # Paso 1: Generar embedding del chunk
    embedding_result = await _generate_embedding_impl(content)
    
    if not embedding_result.get("success"):
        return embedding_result
    
    try:
        # Paso 2: Preparar datos para insertar
        data = {
            "classroom_document_id": classroom_document_id,
            "chunk_index": chunk_index,
            "content": content,
            "embedding": embedding_result["embedding"]
        }
        
        if token_count is not None:
            data["token"] = token_count
        
        # Paso 3: Insertar en Supabase
        result = await asyncio.to_thread(
            lambda: supabase_client.client.table("classroom_document_chunks").insert(data).execute()
        )
        
        if not result.data:
            raise Exception("No se recibieron datos de Supabase después de insertar")
        
        chunk_id = result.data[0]['id']
        
        return {
            "success": True,
            "message": "Chunk almacenado exitosamente",
            "chunk_id": chunk_id,
            "classroom_document_id": classroom_document_id,
            "chunk_index": chunk_index,
            "embedding_dimension": embedding_result["dimension"],
            "content_length": len(content)
        }
    
    except Exception as e:
        error_details = str(e)
        return {
            "success": False,
            "error": f"Error almacenando chunk: {error_details}"
        }


# ====== FUNCIÓN AUXILIAR PARA IMPLEMENTACIÓN ======

async def _store_document_chunks_impl(
    classroom_document_id: str,
    chunk_size: int = 1000,
    chunk_overlap: int = 100
) -> Dict[str, Any]:
    """
    Implementación interna para procesar y almacenar chunks de un documento.
    
    Esta función:
    1. Lee el documento de Supabase Storage
    2. Extrae el contenido (con OCR si es imagen)
    3. Divide en chunks con overlap
    4. Genera embeddings para cada chunk
    5. Almacena los chunks en la base de datos
    
    Args:
        classroom_document_id: UUID del documento en classroom_documents
        chunk_size: Tamaño de cada chunk en caracteres (default: 1000)
        chunk_overlap: Overlap entre chunks en caracteres (default: 100)
        
    Returns:
        Dict con el resultado del procesamiento y chunks creados
    """
    print(f"\n{'='*60}")
    print("📦 IMPL: _store_document_chunks_impl")
    print(f"{'='*60}")
    print(f"📥 Parámetros:")
    print(f"   - classroom_document_id: {classroom_document_id}")
    print(f"   - chunk_size: {chunk_size}")
    print(f"   - chunk_overlap: {chunk_overlap}")
    
    try:
        # Paso 1: Obtener información del documento desde classroom_documents
        print("   🔄 PASO 1: Obteniendo información del documento...")
        
        doc_result = await asyncio.to_thread(
            lambda: supabase_client.client.table("classroom_documents")
            .select("*")
            .eq("id", classroom_document_id)
            .single()
            .execute()
        )
        
        if not doc_result.data:
            raise Exception(f"Documento {classroom_document_id} no encontrado")
        
        doc = doc_result.data
        storage_path = doc.get('storage_path')
        bucket = doc.get('bucket', 'uploads')
        mime_type = doc.get('mime_type', '')
        
        print(f"   ✅ Documento encontrado: {doc.get('title', 'Sin título')}")
        print(f"   📁 Ruta: {storage_path}")
        print(f"   📦 Bucket: {bucket}")
        print(f"   📄 Tipo: {mime_type}")
        
        # Paso 2: Determinar el tipo de documento y extraer contenido
        image_extensions = ['.jpg', '.jpeg', '.png', '.gif', '.webp', '.bmp', '.heic', '.heif']
        image_mime_types = ['image/jpeg', 'image/png', 'image/gif', 'image/webp', 'image/bmp']
        pdf_extensions = ['.pdf']
        pdf_mime_types = ['application/pdf']
        
        is_image = (any(storage_path.lower().endswith(ext) for ext in image_extensions) or 
                    mime_type in image_mime_types)
        is_pdf = (any(storage_path.lower().endswith(ext) for ext in pdf_extensions) or 
                  mime_type in pdf_mime_types)
        
        content = ""
        
        if is_image:
            # Paso 2a: Procesar imagen con OCR
            print(f"   🖼️  Detectada IMAGEN - Aplicando OCR...")
            ocr_result = await _extract_text_from_image_impl(storage_path, bucket)
            
            if not ocr_result.get("success"):
                raise Exception(f"Error en OCR: {ocr_result.get('error')}")
            
            content = ocr_result.get("extracted_text", "")
            print(f"   ✅ Texto extraído por OCR ({len(content)} caracteres)")
            
        elif is_pdf:
            # Paso 2b: Procesar PDF con PyPDF2
            print(f"   📄 Detectado PDF - Extrayendo texto...")
            file_data = await asyncio.to_thread(
                lambda: supabase_client.client.storage.from_(bucket).download(storage_path)
            )
            
            if not file_data:
                raise Exception("No se pudo descargar el archivo PDF")
            
            try:
                # Crear un objeto BytesIO para PyPDF2
                pdf_file = io.BytesIO(file_data)
                pdf_reader = PdfReader(pdf_file)
                
                # Extraer texto de todas las páginas
                text_parts = []
                for page_num, page in enumerate(pdf_reader.pages, 1):
                    page_text = page.extract_text()
                    if page_text.strip():
                        text_parts.append(page_text)
                
                content = "\n\n".join(text_parts)
                print(f"   ✅ Texto extraído del PDF ({len(pdf_reader.pages)} páginas, {len(content)} caracteres)")
                
                if not content.strip():
                    print(f"   ⚠️  El PDF no tiene texto extraíble, podría ser escaneado")
                    print(f"   🖼️  Intentando OCR con Gemini...")
                    # Si el PDF no tiene texto, intentar OCR
                    ocr_result = await _extract_text_from_image_impl(storage_path, bucket)
                    if ocr_result.get("success"):
                        content = ocr_result.get("extracted_text", "")
                        print(f"   ✅ Texto extraído por OCR ({len(content)} caracteres)")
                    
            except Exception as pdf_error:
                print(f"   ❌ Error al procesar PDF: {pdf_error}")
                raise Exception(f"Error al extraer texto del PDF: {pdf_error}")
            
        else:
            # Paso 2c: Descargar y leer archivo de texto plano
            print(f"   📄 Detectado TEXTO PLANO - Descargando...")
            file_data = await asyncio.to_thread(
                lambda: supabase_client.client.storage.from_(bucket).download(storage_path)
            )
            
            if not file_data:
                raise Exception("No se pudo descargar el archivo")
            
            # Intentar decodificar como texto
            try:
                content = file_data.decode('utf-8')
            except UnicodeDecodeError:
                try:
                    content = file_data.decode('latin-1')
                except:
                    raise Exception("No se pudo decodificar el archivo como texto")
            
            print(f"   ✅ Texto descargado ({len(content)} caracteres)")
        
        if not content or len(content.strip()) < 10:
            raise Exception("El contenido extraído es demasiado corto o vacío")
        
        # Paso 2d: Limpiar formato del texto (eliminar saltos de línea excesivos)
        print(f"   🧹 Limpiando formato del texto...")
        original_length = len(content)
        
        # Normalizar espacios en blanco: múltiples espacios/saltos → un espacio
        import re
        content = re.sub(r'\s+', ' ', content)
        content = content.strip()
        
        cleaned_length = len(content)
        reduction = 100 - (cleaned_length / original_length * 100) if original_length > 0 else 0
        print(f"   ✅ Formato limpio ({cleaned_length} caracteres, -{reduction:.1f}%)")
        
        # Paso 3: Dividir en chunks con overlap
        print(f"   🔄 PASO 3: Dividiendo en chunks (size={chunk_size}, overlap={chunk_overlap})...")
        chunks = []
        
        for i in range(0, len(content), chunk_size - chunk_overlap):
            chunk_text = content[i:i + chunk_size]
            if chunk_text.strip():
                chunks.append({
                    'index': len(chunks),
                    'content': chunk_text,
                    'start_pos': i,
                    'end_pos': min(i + chunk_size, len(content))
                })
        
        print(f"   ✅ Creados {len(chunks)} chunks")
        
        # Paso 4: Almacenar cada chunk
        print(f"   🔄 PASO 4: Almacenando {len(chunks)} chunks...")
        stored_chunks = []
        
        for chunk in chunks:
            # Usar la implementación interna directamente (sin logs del tool)
            chunk_result = await _store_document_chunk_impl(
                classroom_document_id=classroom_document_id,
                chunk_index=chunk['index'],
                content=chunk['content'],
                token_count=len(chunk['content'].split())
            )
            
            if chunk_result.get("success"):
                stored_chunks.append({
                    "chunk_id": chunk_result.get("chunk_id"),
                    "chunk_index": chunk['index'],
                    "content_length": len(chunk['content'])
                })
                print(f"   ✅ Chunk {chunk['index']} almacenado")
            else:
                print(f"   ⚠️  Error en chunk {chunk['index']}: {chunk_result.get('error')}")
        
        print(f"✅ Proceso completado exitosamente")
        print(f"   📊 Total chunks: {len(stored_chunks)}/{len(chunks)}")
        print(f"   📝 Total caracteres: {len(content)}")
        print(f"{'='*60}\n")
        
        return {
            "success": True,
            "message": "Documento procesado y almacenado exitosamente",
            "document_id": classroom_document_id,
            "document_length": len(content),
            "chunk_size": chunk_size,
            "chunk_overlap": chunk_overlap,
            "total_chunks": len(stored_chunks),
            "chunks": stored_chunks
        }
        
    except Exception as e:
        error_details = str(e)
        print(f"\n❌ ERROR procesando documento: {error_details}")
        
        return {
            "success": False,
            "error": f"Error procesando documento: {error_details}"
        }


# ====== TOOL: store_document_chunks (versión automática) ======

@mcp.tool()
async def store_document_chunks(
    classroom_document_id: str,
    chunk_size: int = 1000,
    chunk_overlap: int = 100
) -> Dict[str, Any]:
    """
    Procesa y almacena chunks de un documento AUTOMÁTICAMENTE.
    
    Esta es la versión SIMPLIFICADA - solo pasa el ID del documento y listo!
    
    La función hace TODO por ti:
    - Lee el documento de Storage (texto o imagen)
    - Aplica OCR si es imagen
    - Divide en chunks con overlap
    - Genera embeddings
    - Almacena en la base de datos
    
    Args:
        classroom_document_id: UUID del documento en classroom_documents
        chunk_size: Tamaño de cada chunk en caracteres (default: 1000)
        chunk_overlap: Overlap entre chunks en caracteres (default: 100)
        
    Returns:
        Dict con el resultado del procesamiento y chunks creados
    """
    return await _store_document_chunks_impl(
        classroom_document_id=classroom_document_id,
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap
    )


@mcp.tool()
async def search_similar_chunks(
    query_text: str,
    classroom_id: str,
    limit: int = 5,
    threshold: Optional[float] = None
) -> Dict[str, Any]:
    """
    Busca chunks/fragmentos de documentos similares usando búsqueda semántica por embeddings.
    
    Genera un embedding del query y busca los chunks más similares
    SOLO dentro del classroom especificado usando distancia coseno.
    
    Args:
        query_text: Texto de consulta para buscar chunks similares
        classroom_id: UUID del classroom para filtrar (OBLIGATORIO)
        limit: Número máximo de resultados (default: 5)
        threshold: Umbral mínimo de similitud 0-1 (default: 0.6 desde config)
        
    Returns:
        Dict con los chunks similares encontrados y metadata
    """
    print(f"\n{'='*60}")
    print("🎯 TOOL: search_similar_chunks")
    print(f"{'='*60}")
    print(f"📥 Parámetros:")
    print(f"   - Query: '{query_text[:50]}...'")
    print(f"   - classroom_id: {classroom_id}")
    print(f"   - limit: {limit}")
    print(f"   - threshold: {threshold or config.SIMILARITY_THRESHOLD}")
    
    # Usar threshold de config si no se proporciona
    if threshold is None:
        threshold = config.SIMILARITY_THRESHOLD
    
    print(f"\n🔍 Iniciando búsqueda en classroom {classroom_id}...")
    
    # Paso 1: Generar embedding del query
    print("   🔄 PASO 1: Generando embedding del query...")
    embedding_result = await generate_embedding(query_text)
    
    if not embedding_result.get("success"):
        print(f"   ❌ Fallo al generar embedding")
        return embedding_result
    
    print(f"   ✅ Embedding del query generado ({embedding_result.get('dimension')} dims)")
    
    try:
        # Paso 2: Buscar chunks usando función RPC
        print(f"   🔄 PASO 2: Buscando chunks en Supabase...")
        print(f"   📞 Usando match_classroom_chunks RPC")
        
        # Llamar función RPC de Supabase para búsqueda semántica
        result = await asyncio.to_thread(
            lambda: supabase_client.client.rpc(
                'match_classroom_chunks',
                {
                    'query_embedding': embedding_result["embedding"],
                    'filter_classroom_id': classroom_id,
                    'match_threshold': threshold,
                    'match_count': limit
                }
            ).execute()
        )
        
        print(f"   ✅ RPC ejecutado")
        chunks = result.data if result.data else []
        count = len(chunks)
        
        print(f"✅ Búsqueda completada: {count} chunks encontrados")
        if count > 0:
            print(f"   📄 Chunk IDs: {[chunk.get('id') for chunk in chunks[:3]]}")
            print(f"   📊 Similitudes: {[round(chunk.get('similarity', 0), 3) for chunk in chunks[:3]]}")
        print(f"{'='*60}\n")
        
        return {
            "success": True,
            "query": query_text,
            "classroom_id": classroom_id,
            "results": chunks,
            "count": count,
            "threshold_used": threshold,
            "embedding_dimension": embedding_result["dimension"]
        }
    
    except Exception as e:
        error_details = str(e)
        print(f"\n❌ ERROR en búsqueda: {error_details}")
        
        hint = "Verifica que la función RPC 'match_classroom_chunks' exista en Supabase"
        if "function" in error_details.lower() and "does not exist" in error_details.lower():
            hint = "La función match_classroom_chunks no existe. Debes crearla en Supabase."
        
        print(f"   💡 {hint}")
        print(f"{'='*60}\n")
        
        return {
            "success": False,
            "error": f"Error en búsqueda: {error_details}",
            "hint": hint
        }


async def _chat_with_classroom_assistant_impl(request: ChatRequest) -> Dict[str, Any]:
    """
    Chat con el asistente de EstudIA especializado en los documentos del aula.
    
    Proporciona respuestas contextualizadas basándose en los documentos
    subidos al classroom específico. Similar a NotebookLM.
    
    INCLUYE PERSONALIZACIÓN basada en el contexto del usuario (nivel educativo,
    estilo de aprendizaje, preferencias, etc.)
    
    Args:
        request: Objeto con message, classroom_id, user_id opcional y session_id opcional
    
    Returns:
        Dict con la respuesta del asistente y los chunks referenciados
    """
    try:
        print(f"\n{'='*60}")
        print("💬 Chat con asistente de classroom")
        print(f"{'='*60}")
        print(f"   - Message: {request.message[:50]}...")
        print(f"   - Classroom ID: {request.classroom_id}")
        print(f"   - User ID: {request.user_id or 'Anonymous'}")
        
        # PASO 0: Obtener contexto del usuario si está disponible
        user_context_info = ""
        if request.user_id:
            try:
                print(f"   👤 Obteniendo contexto del usuario...")
                user_result = await asyncio.to_thread(
                    lambda: supabase_client.client.table("users")
                    .select("user_context, name")
                    .eq("id", request.user_id)
                    .single()
                    .execute()
                )
                
                if user_result.data:
                    user_context = user_result.data.get('user_context', '')
                    user_name = user_result.data.get('name', 'Estudiante')
                    
                    if user_context:
                        user_context_info = f"""
**CONTEXTO DEL ESTUDIANTE ({user_name}):**
{user_context}

IMPORTANTE: Adapta tu respuesta según este contexto:
- Usa el nivel de complejidad apropiado para su nivel educativo
- Considera su estilo de aprendizaje preferido
- Ten en cuenta sus fortalezas y áreas de mejora
- Respeta sus preferencias de comunicación
- Personaliza ejemplos según sus intereses
"""
                        print(f"   ✅ Contexto del usuario obtenido ({len(user_context)} caracteres)")
                    else:
                        print(f"   ℹ️  Usuario sin contexto personalizado")
            except Exception as e:
                print(f"   ⚠️  No se pudo obtener contexto del usuario: {e}")
        
        # PASO 1: Generar embedding para encontrar chunks relevantes
        embedding = await gemini_client.generate_embedding(request.message)
        
        # PASO 2: Buscar chunks relevantes en el classroom usando RPC directamente
        try:
            result = await asyncio.to_thread(
                lambda: supabase_client.client.rpc(
                    'match_classroom_chunks',
                    {
                        'query_embedding': embedding,
                        'filter_classroom_id': request.classroom_id,
                        'match_threshold': 0.5,
                        'match_count': 5
                    }
                ).execute()
            )
            relevant_chunks = result.data if result.data else []
        except Exception as e:
            print(f"   ⚠️  Error buscando chunks: {e}")
            relevant_chunks = []
        
        print(f"   📚 Chunks encontrados: {len(relevant_chunks)}")
        
        # PASO 3: Construir contexto con los chunks y extraer IDs de documentos únicos
        context_blocks = []
        document_ids = set()  # Para almacenar IDs únicos de documentos
        documents_info = {}   # Para almacenar información detallada de cada documento
        
        for i, chunk in enumerate(relevant_chunks, start=1):
            content = chunk.get('content', '')
            doc_id = chunk.get('classroom_document_id', 'Unknown')
            chunk_idx = chunk.get('chunk_index', 0)
            similarity = chunk.get('similarity', 0)
            
            # Agregar el ID del documento al set
            if doc_id != 'Unknown':
                document_ids.add(doc_id)
                
                # Guardar información del documento si no existe
                if doc_id not in documents_info:
                    documents_info[doc_id] = {
                        'document_id': doc_id,
                        'chunks_used': [],
                        'max_similarity': similarity
                    }
                
                # Agregar información del chunk usado
                documents_info[doc_id]['chunks_used'].append({
                    'chunk_index': chunk_idx,
                    'similarity': similarity,
                    'content_preview': content[:100] + '...' if len(content) > 100 else content
                })
                
                # Actualizar similitud máxima si es mayor
                if similarity > documents_info[doc_id]['max_similarity']:
                    documents_info[doc_id]['max_similarity'] = similarity
            
            block = f"[Chunk {i} - Doc: {doc_id}, Index: {chunk_idx}, Similitud: {similarity:.3f}]\n{content}"
            context_blocks.append(block)
        
        context = "\n\n---\n\n".join(context_blocks) if context_blocks else "No se encontraron documentos relevantes."
        
        # PASO 3.5: Obtener información adicional de los documentos desde la tabla
        documents_details = []
        if document_ids:
            try:
                print(f"   📄 Obteniendo detalles de {len(document_ids)} documentos...")
                docs_result = await asyncio.to_thread(
                    lambda: supabase_client.client.table("classroom_documents")
                    .select("id, title, description, original_filename, mime_type, storage_path, bucket")
                    .in_("id", list(document_ids))
                    .execute()
                )
                
                if docs_result.data:
                    for doc in docs_result.data:
                        doc_id = doc['id']
                        doc_info = documents_info.get(doc_id, {})
                        
                        documents_details.append({
                            'document_id': doc_id,
                            'title': doc.get('title', 'Sin título'),
                            'description': doc.get('description'),
                            'filename': doc.get('original_filename'),
                            'mime_type': doc.get('mime_type'),
                            'storage_path': doc.get('storage_path'),
                            'bucket': doc.get('bucket', 'uploads'),
                            'chunks_used': doc_info.get('chunks_used', []),
                            'relevance_score': doc_info.get('max_similarity', 0)
                        })
                    
                    # Ordenar por relevancia (mayor similitud primero)
                    documents_details.sort(key=lambda x: x['relevance_score'], reverse=True)
                    print(f"   ✅ Detalles de documentos obtenidos")
            except Exception as e:
                print(f"   ⚠️  Error obteniendo detalles de documentos: {e}")
        
        # PASO 4: Obtener respuesta del asistente con contexto personalizado
        print(f"   🤖 Generando respuesta personalizada con Gemini...")
        
        prompt = f"""Eres un asistente educativo que ayuda a estudiantes respondiendo preguntas basándote en los documentos de su aula.

{user_context_info}

**Pregunta del estudiante:**
{request.message}

**Documentos relevantes del aula:**
{context}

**Instrucciones:**
- Responde basándote ÚNICAMENTE en la información de los documentos proporcionados
- Si la información no está en los documentos, indícalo claramente
- Sé claro, conciso y educativo
- ADAPTA tu lenguaje y complejidad según el contexto del estudiante
- Si el estudiante prefiere aprendizaje visual, menciona diagramas o imágenes cuando sea relevante
- Si prefiere código/práctica, enfócate en ejemplos prácticos
- Ajusta la profundidad de tu explicación según su nivel educativo
- Usa un tono y vocabulario apropiado para su contexto
- Cita específicamente qué chunk/documento usaste si es relevante
- Si no hay documentos relevantes, sugiere reformular la pregunta o subir documentos sobre el tema
"""
        
        # Generar respuesta con Gemini
        response = await gemini_client.generate_text(prompt)
        
        print(f"   ✅ Respuesta personalizada generada")
        print(f"   📚 Documentos únicos referenciados: {len(document_ids)}")
        if documents_details:
            for doc in documents_details[:3]:
                print(f"      - {doc['title']} (relevancia: {doc['relevance_score']:.3f})")
        print(f"{'='*60}\n")
        
        return {
            'success': True,
            'data': {
                'response': response,
                'chunks_referenced': len(relevant_chunks),
                'chunks': relevant_chunks[:3],  # Solo las 3 más relevantes (backward compatibility)
                'classroom_id': request.classroom_id,
                'personalized': bool(user_context_info),  # Indica si se personalizó
                # NUEVO: Información estructurada de documentos para preview
                'documents': documents_details,  # Lista completa con detalles de cada documento
                'document_ids': list(document_ids),  # Solo los IDs únicos
                'total_documents': len(document_ids)
            },
            'message': "Respuesta del asistente generada"
        }
        
    except Exception as error:
        print(f"❌ Error en chat: {error}")
        return {
            'success': False,
            'error': str(error),
            'message': "Error en el chat con el asistente"
        }


@mcp.tool()
async def chat_with_classroom_assistant(request: ChatRequest) -> Dict[str, Any]:
    """
    Chat con el asistente de EstudIA especializado en los documentos del aula.
    
    Proporciona respuestas contextualizadas basándose en los documentos
    subidos al classroom específico. Similar a NotebookLM.
    
    INCLUYE PERSONALIZACIÓN basada en el contexto del usuario (nivel educativo,
    estilo de aprendizaje, preferencias, etc.)
    
    Args:
        request: Objeto con message, classroom_id, user_id opcional y session_id opcional
    
    Returns:
        Dict con la respuesta del asistente y los chunks referenciados
    """
    return await _chat_with_classroom_assistant_impl(request)





async def _analyze_and_update_user_context_impl(
    user_id: str,
    session_id: str
) -> Dict[str, Any]:
    """
    Implementación interna de analyze_and_update_user_context.
    Esta función contiene la lógica real y puede ser llamada directamente.
    """
    print(f"\n{'='*70}")
    print("🎯 TOOL: analyze_and_update_user_context")
    print(f"{'='*70}")
    print(f"📥 Parámetros:")
    print(f"   - User ID: {user_id}")
    print(f"   - Session ID: {session_id}")
    
    try:
        # PASO 1: Obtener el contexto actual del usuario
        print(f"\n👤 PASO 1: Obteniendo contexto actual del usuario...")
        
        user_result = await asyncio.to_thread(
            lambda: supabase_client.client.table("users")
            .select("user_context, name, email")
            .eq("id", user_id)
            .single()
            .execute()
        )
        
        if not user_result.data:
            return {
                "success": False,
                "error": "Usuario no encontrado",
                "user_id": user_id
            }
        
        current_context = user_result.data.get('user_context') or ""
        user_name = user_result.data.get('name', 'Usuario')
        
        print(f"✅ Usuario encontrado: {user_name}")
        print(f"   📝 Contexto actual: {len(current_context)} caracteres")
        
        # PASO 2: Obtener todos los mensajes de la sesión
        print(f"\n💬 PASO 2: Obteniendo mensajes de la sesión...")
        
        messages_result = await asyncio.to_thread(
            lambda: supabase_client.client.table("cubicle_messages")
            .select("id, user_id, content, created_at")
            .eq("session_id", session_id)
            .order("created_at", desc=False)
            .execute()
        )
        
        messages = messages_result.data if messages_result.data else []
        
        print(f"✅ Encontrados {len(messages)} mensajes en la sesión")
        
        if len(messages) == 0:
            return {
                "success": True,
                "message": "No hay mensajes en la sesión para analizar",
                "context_updated": False,
                "user_id": user_id,
                "session_id": session_id
            }
        
        # PASO 3: Analizar la conversación con Gemini
        print(f"\n🤖 PASO 3: Analizando conversación con Gemini...")
        print(f"   📊 Total mensajes a analizar: {len(messages)}")
        
        analysis = await gemini_client.analyze_conversation_for_context_update(
            current_context=current_context,
            conversation_messages=messages
        )
        
        should_update = analysis.get('should_update', False)
        new_context = analysis.get('new_context', current_context)
        reasons = analysis.get('reasons', [])
        key_findings = analysis.get('key_findings', {})
        
        print(f"✅ Análisis completado")
        print(f"   🔄 ¿Debe actualizarse?: {should_update}")
        print(f"   📋 Razones: {len(reasons)}")
        
        # PASO 4: Actualizar el contexto si es necesario
        if should_update:
            print(f"\n💾 PASO 4: Actualizando contexto del usuario...")
            
            update_result = await asyncio.to_thread(
                lambda: supabase_client.client.table("users")
                .update({"user_context": new_context})
                .eq("id", user_id)
                .execute()
            )
            
            print(f"✅ Contexto actualizado exitosamente")
            print(f"   📝 Nuevo contexto: {len(new_context)} caracteres")
            
            for i, reason in enumerate(reasons, 1):
                print(f"   {i}. {reason}")
        else:
            print(f"\n⏭️  PASO 4: No se requiere actualización")
            for reason in reasons:
                print(f"   • {reason}")
        
        print(f"{'='*70}\n")
        
        return {
            "success": True,
            "context_updated": should_update,
            "previous_context": current_context,
            "new_context": new_context if should_update else current_context,
            "reasons": reasons,
            "key_findings": key_findings,
            "messages_analyzed": len(messages),
            "user_id": user_id,
            "user_name": user_name,
            "session_id": session_id,
            "message": f"Contexto {'actualizado' if should_update else 'sin cambios'}"
        }
    
    except Exception as e:
        error_details = str(e)
        print(f"\n❌ ERROR: {error_details}")
        print(f"{'='*70}\n")
        
        return {
            "success": False,
            "error": f"Error analizando contexto: {error_details}",
            "user_id": user_id,
            "session_id": session_id
        }


@mcp.tool()
async def analyze_and_update_user_context(
    user_id: str,
    session_id: str
) -> Dict[str, Any]:
    """
    Analiza la conversación de una sesión de cubículo y actualiza el contexto del usuario.
    
    Esta herramienta lee todos los mensajes de la sesión, los analiza junto con el contexto
    actual del usuario, y determina si debe actualizarse con nueva información relevante
    sobre nivel educativo, estilo de aprendizaje, preferencias, etc.
    
    Args:
        user_id: UUID del usuario cuyo contexto se analizará
        session_id: UUID de la sesión del cubículo con la conversación
        
    Returns:
        Dict con el resultado del análisis y si se actualizó el contexto
    """
    return await _analyze_and_update_user_context_impl(user_id, session_id)


async def _generate_resources_impl(
    classroom_id: str,
    resource_type: str,
    user_id: str,
    topic: str = None,
    source_document_ids: list = None
) -> Dict[str, Any]:
    """
    Implementación interna de generate_resources.
    Genera recursos educativos (PDF o PPT) basándose en documentos del classroom.
    """
    import io
    import uuid
    from datetime import datetime
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, PageBreak
    from reportlab.lib.units import inch
    from pptx import Presentation
    from pptx.util import Inches, Pt
    
    print(f"\n{'='*70}")
    print("🎯 TOOL: generate_resources")
    print(f"{'='*70}")
    print(f"📥 Parámetros:")
    print(f"   - Classroom ID: {classroom_id}")
    print(f"   - Resource Type: {resource_type}")
    print(f"   - User ID: {user_id}")
    print(f"   - Topic: {topic or 'General'}")
    
    try:
        # PASO 1: Validar tipo de recurso
        if resource_type not in ['pdf', 'ppt']:
            return {
                "success": False,
                "error": f"Tipo de recurso inválido: {resource_type}. Use 'pdf' o 'ppt'"
            }
        
        # PASO 2: Obtener documentos del classroom
        print(f"\n📚 PASO 2: Obteniendo documentos del classroom...")
        
        # Si se proporcionaron IDs específicos, usarlos (validando que pertenezcan al classroom)
        if source_document_ids:
            print(f"   🔍 Filtrando por {len(source_document_ids)} documentos específicos")
            docs_result = await asyncio.to_thread(
                lambda: supabase_client.client.table("classroom_documents")
                .select("id, title, original_filename, storage_path")
                .eq("classroom_id", classroom_id)  # IMPORTANTE: Validar que pertenezcan al classroom
                .in_("id", source_document_ids)
                .execute()
            )
        else:
            # Obtener todos los documentos del classroom
            docs_result = await asyncio.to_thread(
                lambda: supabase_client.client.table("classroom_documents")
                .select("id, title, original_filename, storage_path")
                .eq("classroom_id", classroom_id)
                .execute()
            )
        
        documents = docs_result.data if docs_result.data else []
        print(f"✅ Encontrados {len(documents)} documentos")
        
        # Validar si se pidieron documentos específicos pero no se encontraron
        if source_document_ids and len(documents) == 0:
            return {
                "success": False,
                "error": "Los documentos especificados no existen o no pertenecen a este classroom"
            }
        elif source_document_ids and len(documents) < len(source_document_ids):
            print(f"   ⚠️  Advertencia: Solo {len(documents)} de {len(source_document_ids)} documentos encontrados")
        
        if not documents:
            return {
                "success": False,
                "error": "No hay documentos disponibles en el classroom"
            }
        
        # Obtener chunks de los documentos
        doc_ids = [doc['id'] for doc in documents]
        chunks_result = await asyncio.to_thread(
            lambda: supabase_client.client.table("classroom_document_chunks")
            .select("content, chunk_index, classroom_document_id")
            .in_("classroom_document_id", doc_ids)
            .order("classroom_document_id")
            .order("chunk_index")
            .limit(100)
            .execute()
        )
        
        chunks = chunks_result.data if chunks_result.data else []
        print(f"✅ Encontrados {len(chunks)} chunks")
        
        # PASO 3: Obtener contexto del usuario para personalización
        print(f"\n👤 PASO 3: Obteniendo contexto del usuario...")
        user_context_info = ""
        user_name = "Estudiante"
        
        try:
            user_result = await asyncio.to_thread(
                lambda: supabase_client.client.table("users")
                .select("user_context, name")
                .eq("id", user_id)
                .single()
                .execute()
            )
            
            if user_result.data:
                user_context = user_result.data.get('user_context', '')
                user_name = user_result.data.get('name', 'Estudiante')
                
                if user_context:
                    user_context_info = f"""
**CONTEXTO DEL ESTUDIANTE ({user_name}):**
{user_context}

IMPORTANTE: Personaliza el recurso según este contexto:
- Ajusta el nivel de complejidad según su nivel educativo
- Considera su estilo de aprendizaje preferido (visual, práctico, teórico, etc.)
- Ten en cuenta sus fortalezas y áreas de mejora
- Usa vocabulario y ejemplos apropiados para su nivel
- Si prefiere aprendizaje visual, enfatiza diagramas y estructuras visuales
- Si prefiere código/práctica, incluye ejemplos prácticos y aplicaciones
"""
                    print(f"   ✅ Contexto del usuario obtenido ({len(user_context)} caracteres)")
                else:
                    print(f"   ℹ️  Usuario sin contexto personalizado")
        except Exception as e:
            print(f"   ⚠️  No se pudo obtener contexto del usuario: {e}")
        
        # PASO 4: Preparar contenido para Gemini
        print(f"\n📝 PASO 4: Preparando contenido...")
        
        full_content = "\n\n".join([
            chunk.get('content', '') for chunk in chunks[:30]
        ])
        
        print(f"✅ Contenido preparado ({len(full_content)} caracteres)")
        
        # PASO 5: Generar estructura con Gemini (PERSONALIZADA)
        print(f"\n🤖 PASO 5: Generando estructura personalizada del recurso con Gemini...")
        
        topic_text = f" sobre '{topic}'" if topic else ""
        
        prompt = f"""Eres un asistente pedagógico experto. Genera un recurso educativo{topic_text} basado en el siguiente contenido.

{user_context_info}

**Contenido de los documentos:**
{full_content}

ADAPTA el contenido del recurso según el contexto del estudiante si está disponible.

**Contenido de los documentos:**
{full_content}

Genera una estructura en formato JSON con:

{{
  "title": "Título del recurso",
  "subtitle": "Subtítulo o descripción breve",
  "sections": [
    {{
      "heading": "Título de la sección",
      "content": "Contenido detallado de la sección (2-3 párrafos)",
      "bullet_points": ["Punto clave 1", "Punto clave 2", "Punto clave 3"]
    }}
  ],
  "key_concepts": [
    {{
      "term": "Término",
      "definition": "Definición clara y concisa"
    }}
  ],
  "summary": "Resumen final del recurso (1-2 párrafos)"
}}

Responde SOLO con JSON válido:"""
        
        response = await gemini_client.generate_text(prompt)
        
        # Parsear JSON
        try:
            json_str = response.strip()
            if json_str.startswith("```json"):
                json_str = json_str[7:]
            if json_str.startswith("```"):
                json_str = json_str[3:]
            if json_str.endswith("```"):
                json_str = json_str[:-3]
            
            structure = json.loads(json_str.strip())
            print(f"✅ Estructura personalizada generada: {len(structure.get('sections', []))} secciones")
            
        except json.JSONDecodeError as e:
            return {
                "success": False,
                "error": f"Error parseando estructura: {str(e)}"
            }
        
        # PASO 6: Generar archivo según el tipo
        print(f"\n📄 PASO 6: Generando archivo {resource_type.upper()}...")
        
        file_buffer = io.BytesIO()
        
        if resource_type == 'pdf':
            # Generar PDF
            doc = SimpleDocTemplate(file_buffer, pagesize=letter)
            story = []
            styles = getSampleStyleSheet()
            
            # Título
            title_style = styles['Title']
            story.append(Paragraph(structure.get('title', 'Recurso Educativo'), title_style))
            story.append(Spacer(1, 0.2*inch))
            
            # Subtítulo
            if structure.get('subtitle'):
                story.append(Paragraph(structure['subtitle'], styles['Heading2']))
                story.append(Spacer(1, 0.3*inch))
            
            # Secciones
            for section in structure.get('sections', []):
                story.append(Paragraph(section.get('heading', 'Sección'), styles['Heading1']))
                story.append(Spacer(1, 0.1*inch))
                story.append(Paragraph(section.get('content', ''), styles['BodyText']))
                story.append(Spacer(1, 0.2*inch))
                
                # Bullet points
                for point in section.get('bullet_points', []):
                    story.append(Paragraph(f"• {point}", styles['BodyText']))
                
                story.append(Spacer(1, 0.3*inch))
            
            # Conceptos clave
            if structure.get('key_concepts'):
                story.append(PageBreak())
                story.append(Paragraph("Conceptos Clave", styles['Heading1']))
                story.append(Spacer(1, 0.2*inch))
                
                for concept in structure['key_concepts']:
                    story.append(Paragraph(f"<b>{concept.get('term', '')}</b>: {concept.get('definition', '')}", styles['BodyText']))
                    story.append(Spacer(1, 0.1*inch))
            
            # Resumen
            if structure.get('summary'):
                story.append(PageBreak())
                story.append(Paragraph("Resumen", styles['Heading1']))
                story.append(Spacer(1, 0.2*inch))
                story.append(Paragraph(structure['summary'], styles['BodyText']))
            
            doc.build(story)
            
        else:  # ppt
            # Generar PowerPoint
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Slide 1: Título
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = structure.get('title', 'Recurso Educativo')
            subtitle.text = structure.get('subtitle', '')
            
            # Slides para cada sección
            for section in structure.get('sections', []):
                slide_layout = prs.slide_layouts[1]  # Title and Content
                slide = prs.slides.add_slide(slide_layout)
                
                title = slide.shapes.title
                title.text = section.get('heading', 'Sección')
                
                # Contenido
                content_shape = slide.placeholders[1]
                text_frame = content_shape.text_frame
                text_frame.clear()
                
                # Agregar contenido como párrafo
                p = text_frame.paragraphs[0]
                p.text = section.get('content', '')[:300] + "..."  # Limitar texto
                p.font.size = Pt(14)
                p.level = 0
                
                # Agregar bullet points
                for point in section.get('bullet_points', [])[:5]:  # Max 5 bullets
                    p = text_frame.add_paragraph()
                    p.text = point
                    p.font.size = Pt(12)
                    p.level = 1
            
            # Slide: Conceptos clave
            if structure.get('key_concepts'):
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                title.text = "Conceptos Clave"
                
                content_shape = slide.placeholders[1]
                text_frame = content_shape.text_frame
                text_frame.clear()
                
                for concept in structure['key_concepts'][:6]:  # Max 6 conceptos
                    p = text_frame.add_paragraph() if text_frame.paragraphs[0].text else text_frame.paragraphs[0]
                    p.text = f"{concept.get('term', '')}: {concept.get('definition', '')}"
                    p.font.size = Pt(12)
                    p.level = 0
            
            # Slide final: Resumen
            if structure.get('summary'):
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                title.text = "Resumen"
                
                content_shape = slide.placeholders[1]
                text_frame = content_shape.text_frame
                text_frame.clear()
                
                p = text_frame.paragraphs[0]
                p.text = structure['summary']
                p.font.size = Pt(14)
            
            prs.save(file_buffer)
        
        file_buffer.seek(0)
        file_data = file_buffer.read()
        file_size = len(file_data)
        
        print(f"✅ Archivo generado: {file_size} bytes")
        
        # PASO 7: Subir a Supabase Storage
        print(f"\n☁️  PASO 7: Subiendo archivo a Supabase Storage...")
        
        resource_id = str(uuid.uuid4())
        file_extension = 'pdf' if resource_type == 'pdf' else 'pptx'
        filename = f"{resource_id}.{file_extension}"
        storage_path = f"{classroom_id}/{user_id}/{filename}"
        
        mime_type = 'application/pdf' if resource_type == 'pdf' else 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        
        # Subir archivo (usar 'uploads' si 'generated-resources' no existe)
        bucket_name = 'uploads'  # Cambiar a 'generated-resources' cuando el bucket exista
        upload_result = await asyncio.to_thread(
            lambda: supabase_client.client.storage.from_(bucket_name).upload(
                path=storage_path,
                file=file_data,
                file_options={
                    "content-type": mime_type,
                    "upsert": "false"
                }
            )
        )
        
        print(f"✅ Archivo subido: {storage_path}")
        
        # PASO 8: Guardar metadata en la base de datos
        print(f"\n💾 PASO 8: Guardando metadata en la base de datos...")
        
        resource_data = {
            "id": resource_id,
            "classroom_id": classroom_id,
            "user_id": user_id,
            "resource_type": resource_type,
            "title": structure.get('title', 'Recurso Educativo'),
            "bucket": bucket_name,  # Usar el bucket que configuramos arriba
            "storage_path": storage_path,
            "file_size_bytes": file_size,
            "mime_type": mime_type,
            "generated_with_model": "gemini-2.0-flash",
            "generation_prompt": topic or "Recurso general del classroom",
            "source_document_ids": doc_ids,
            "created_at": datetime.utcnow().isoformat()
        }
        
        insert_result = await asyncio.to_thread(
            lambda: supabase_client.client.table("generated_resources")
            .insert(resource_data)
            .execute()
        )
        
        print(f"✅ Metadata guardada con ID: {resource_id}")
        
        # Obtener URL pública
        public_url = supabase_client.client.storage.from_(bucket_name).get_public_url(storage_path)
        
        personalized = bool(user_context_info)
        if personalized:
            print(f"✨ Recurso personalizado para: {user_name}")
        
        print(f"{'='*70}\n")
        
        return {
            "success": True,
            "message": f"Recurso {resource_type.upper()} {'personalizado' if personalized else 'generado'} exitosamente",
            "resource_id": resource_id,
            "resource_type": resource_type,
            "title": structure.get('title'),
            "storage_path": storage_path,
            "bucket": bucket_name,
            "personalized": personalized,
            "user_name": user_name,
            "file_size_bytes": file_size,
            "public_url": public_url,
            "sections_count": len(structure.get('sections', [])),
            "concepts_count": len(structure.get('key_concepts', [])),
            "source_documents": len(doc_ids)
        }
    
    except Exception as e:
        error_details = str(e)
        print(f"\n❌ ERROR: {error_details}")
        print(f"{'='*70}\n")
        
        return {
            "success": False,
            "error": f"Error generando recurso: {error_details}"
        }


@mcp.tool()
async def generate_resources(
    classroom_id: str,
    resource_type: str,
    user_id: str,
    topic: str = None,
    source_document_ids: list = None
) -> Dict[str, Any]:
    """
    Genera recursos educativos (PDF o PowerPoint) basándose en documentos del classroom.
    
    El recurso se sube automáticamente a Supabase Storage y devuelve la URL pública
    para descargar desde la aplicación React Native.
    
    Args:
        classroom_id: UUID del classroom
        resource_type: Tipo de recurso ('pdf' o 'ppt')
        user_id: UUID del usuario que solicita el recurso
        topic: (Opcional) Tema específico para el recurso
        source_document_ids: (Opcional) Lista de IDs de documentos específicos a usar
        
    Returns:
        Dict con información del recurso generado y URL de descarga
    """
    return await _generate_resources_impl(
        classroom_id=classroom_id,
        resource_type=resource_type,
        user_id=user_id,
        topic=topic,
        source_document_ids=source_document_ids
    )



# ====== FUNCIÓN PRINCIPAL ======

def main():
    """Función principal para ejecutar el servidor MCP"""
    try:
        print("🚀 Iniciando EstudIA MCP Server con FastMCP...")
        print("📋 Herramientas registradas:")
        print("   ✅ generate_embedding - Generar embeddings de texto")
        print("   ✅ store_document_chunks - Almacenar chunks con embeddings")
        print("   ✅ search_similar_chunks - Buscar chunks similares en classroom")
        print("   ✅ chat_with_classroom_assistant - Chat con asistente del aula")
        print("   ✅ analyze_and_update_user_context - Analizar conversación y actualizar contexto de usuario")
        print("   ✅ generate_resources - Generar recursos educativos (PDF/PPT)")
        print("🎯 Servidor MCP listo para recibir peticiones...")
        
        # Ejecutar el servidor FastMCP
        mcp.run()
        
    except Exception as error:
        print(f"❌ Error iniciando el servidor MCP: {error}")
        raise error

if __name__ == "__main__":
    main()
